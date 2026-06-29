#!/usr/bin/env python3
"""化工行业 PPT 渲染器：JSON 大纲 -> 可编辑 .pptx。

用法：
    python scripts/build_pptx.py --outline outline.json --out 输出.pptx
    python scripts/build_pptx.py --outline outline.json            # 自动按主题命名

大纲 schema 见 references/outline-schema.md。模板色系见 templates/themes.json。
所有页面均用 python-pptx 程序化绘制，输出为标准 OOXML，可在 PowerPoint/WPS 二次编辑。
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

EMU = 914400
SLIDE_W = int(13.333 * EMU)
SLIDE_H = int(7.5 * EMU)
FONT = "微软雅黑"

ROOT = Path(__file__).resolve().parent.parent
THEMES = json.loads((ROOT / "templates" / "themes.json").read_text(encoding="utf-8"))


def hex2rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h.lstrip("#"))


class Deck:
    def __init__(self, theme: dict, base: str | None = None, logo: str | None = None):
        self.t = theme
        self.logo = logo
        if base:  # 以现有 PPT 为母版，继承主题色/字体/版式
            self.prs = Presentation(base)
            self._strip(self.prs)
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_W
            self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]

    @staticmethod
    def _strip(prs):
        lst = prs.slides._sldIdLst
        for sld in list(lst):
            rid = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            if rid:
                prs.part.drop_rel(rid)
            lst.remove(sld)

    # --- low level helpers ---
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _logo(self, s, page, x, y, h):
        path = page.get("logo") or self.logo
        if path and Path(path).exists():
            s.shapes.add_picture(path, x, y, height=h)

    def _box(self, s, x, y, w, h, fill=None, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.shadow.inherit = False
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = hex2rgb(fill)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = hex2rgb(line)
        return shp

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, (txt, size, color, bold) in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = hex2rgb(color)
        return tb

    def _bar(self, s):  # 顶部主色条
        self._box(s, 0, 0, SLIDE_W, int(0.18 * EMU), fill=self.t["primary"])

    def _footer(self, s, n):
        self._text(s, int(0.5 * EMU), int(7.0 * EMU), int(12.3 * EMU), int(0.4 * EMU),
                   [(self.t.get("footer", "中国石油化工集团有限公司"), 9, "999999", False)])
        self._text(s, int(12.3 * EMU), int(7.0 * EMU), int(0.7 * EMU), int(0.4 * EMU),
                   [(str(n), 9, "999999", False)], align=PP_ALIGN.RIGHT)

    # --- page types ---
    def cover(self, p):
        s = self._slide()
        self._box(s, 0, 0, SLIDE_W, SLIDE_H, fill=self.t["primary"])
        self._box(s, 0, int(5.2 * EMU), SLIDE_W, int(0.06 * EMU), fill="FFFFFF")
        self._text(s, int(1.0 * EMU), int(2.4 * EMU), int(11.3 * EMU), int(2.0 * EMU),
                   [(p["title"], 40, "FFFFFF", True)], anchor=MSO_ANCHOR.MIDDLE)
        if p.get("subtitle"):
            self._text(s, int(1.0 * EMU), int(4.4 * EMU), int(11.3 * EMU), int(0.8 * EMU),
                       [(p["subtitle"], 18, "F2D0D0", False)])
        self._text(s, int(1.0 * EMU), int(6.4 * EMU), int(11.3 * EMU), int(0.6 * EMU),
                   [(p.get("org", "中国石油化工集团有限公司"), 14, "FFFFFF", False)])
        self._logo(s, p, int(11.3 * EMU), int(0.6 * EMU), int(0.9 * EMU))

    def toc(self, p):
        s = self._slide(); self._bar(s)
        self._text(s, int(0.8 * EMU), int(0.6 * EMU), int(8 * EMU), int(0.9 * EMU),
                   [("目  录 / CONTENTS", 28, self.t["primary"], True)])
        items = p["items"]
        for i, it in enumerate(items):
            y = int((1.9 + i * 0.95) * EMU)
            self._text(s, int(1.0 * EMU), y, int(0.9 * EMU), int(0.7 * EMU),
                       [(f"{i + 1:02d}", 26, self.t["accent"], True)])
            self._text(s, int(2.0 * EMU), y, int(10 * EMU), int(0.7 * EMU),
                       [(it, 18, "333333", False)], anchor=MSO_ANCHOR.MIDDLE)

    def section(self, p):
        s = self._slide()
        self._box(s, 0, 0, SLIDE_W, SLIDE_H, fill=self.t["primary"])
        self._box(s, int(0.8 * EMU), int(3.1 * EMU), int(1.2 * EMU), int(0.12 * EMU), fill="FFFFFF")
        self._text(s, int(0.8 * EMU), int(3.4 * EMU), int(11.5 * EMU), int(1.2 * EMU),
                   [(p["title"], 36, "FFFFFF", True)])
        if p.get("subtitle"):
            self._text(s, int(0.8 * EMU), int(4.6 * EMU), int(11.5 * EMU), int(0.6 * EMU),
                       [(p["subtitle"], 16, "F2D0D0", False)])

    def points(self, p):
        s = self._slide(); self._bar(s); self._title(s, p["title"])
        items = p.get("points", [])
        for i, it in enumerate(items):
            y = int((1.9 + i * 0.95) * EMU)
            self._box(s, int(0.9 * EMU), y + int(0.08 * EMU), int(0.18 * EMU), int(0.55 * EMU), fill=self.t["accent"])
            head = it.get("head", "") if isinstance(it, dict) else it
            body = it.get("body", "") if isinstance(it, dict) else ""
            runs = [(head, 17, "1A1A1A", True)]
            if body:
                runs.append((body, 13, "666666", False))
            self._text(s, int(1.3 * EMU), y, int(11 * EMU), int(0.85 * EMU), runs)

    def table(self, p):
        s = self._slide(); self._bar(s); self._title(s, p["title"])
        rows = p["rows"]; ncol = len(rows[0])
        tb = s.shapes.add_table(len(rows), ncol, int(0.8 * EMU), int(1.9 * EMU),
                                int(11.7 * EMU), int(4.6 * EMU)).table
        for ci, c in enumerate(rows[0]):
            cell = tb.cell(0, ci); cell.fill.solid(); cell.fill.fore_color.rgb = hex2rgb(self.t["primary"])
            self._cell(cell, c, 13, "FFFFFF", True)
        for ri in range(1, len(rows)):
            for ci, c in enumerate(rows[ri]):
                cell = tb.cell(ri, ci); self._cell(cell, str(c), 12, "333333", False)
        self._footer(s, p["_n"])

    def warning(self, p):
        s = self._slide(); self._bar(s); self._title(s, p["title"])
        self._box(s, int(0.8 * EMU), int(1.9 * EMU), int(11.7 * EMU), int(4.5 * EMU),
                  fill="FCEDED", line=self.t["primary"])
        self._text(s, int(1.1 * EMU), int(2.1 * EMU), int(11 * EMU), int(0.7 * EMU),
                   [("⚠ 安全提示", 20, self.t["primary"], True)])
        for i, it in enumerate(p.get("points", [])):
            self._text(s, int(1.3 * EMU), int((2.9 + i * 0.7) * EMU), int(10.8 * EMU), int(0.6 * EMU),
                       [(f"· {it}", 15, "1A1A1A", False)])
        self._footer(s, p["_n"])

    def closing(self, p):
        s = self._slide()
        self._box(s, 0, 0, SLIDE_W, SLIDE_H, fill=self.t["primary"])
        self._text(s, 0, int(2.8 * EMU), SLIDE_W, int(1.5 * EMU),
                   [(p.get("title", "谢 谢"), 48, "FFFFFF", True)], align=PP_ALIGN.CENTER)
        self._text(s, 0, int(4.4 * EMU), SLIDE_W, int(0.6 * EMU),
                   [(p.get("org", "中国石油化工集团有限公司"), 16, "F2D0D0", False)], align=PP_ALIGN.CENTER)

    def _title(self, s, t):
        self._text(s, int(0.8 * EMU), int(0.55 * EMU), int(11.5 * EMU), int(1.0 * EMU),
                   [(t, 26, self.t["primary"], True)])
        self._box(s, int(0.85 * EMU), int(1.55 * EMU), int(2.0 * EMU), int(0.06 * EMU), fill=self.t["accent"])

    def _cell(self, cell, text, size, color, bold):
        cell.text = text
        for pr in cell.text_frame.paragraphs:
            for r in pr.runs:
                r.font.size = Pt(size); r.font.name = FONT; r.font.bold = bold
                r.font.color.rgb = hex2rgb(color)

    def render(self, pages):
        types = {"cover": self.cover, "toc": self.toc, "section": self.section,
                 "points": self.points, "table": self.table, "warning": self.warning,
                 "closing": self.closing}
        for n, pg in enumerate(pages, 1):
            pg["_n"] = n
            fn = types.get(pg["type"], self.points)
            fn(pg)
            if pg["type"] in ("toc", "points", "section"):
                self._footer(self.prs.slides[-1], n)

    def save(self, path):
        self.prs.save(path)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--outline", required=True)
    ap.add_argument("--out", default=None)
    ap.add_argument("--base", default=None, help="现有 PPT 母版路径，套用其版式")
    ap.add_argument("--logo", default=None, help="logo 图片路径，封面插入")
    a = ap.parse_args()
    data = json.loads(Path(a.outline).read_text(encoding="utf-8"))
    theme = THEMES.get(data.get("template", "项目汇报"), THEMES["项目汇报"])
    base = a.base or data.get("base")
    logo = a.logo or data.get("logo") or theme.get("logo")
    out = a.out or f"{data.get('title', 'output')}.pptx"
    deck = Deck(theme, base=base, logo=logo)
    deck.render(data["pages"])
    deck.save(out)
    print(f"已生成 {out}（{len(data['pages'])} 页，模板：{data.get('template')}，母版：{base or '内置'}）")


if __name__ == "__main__":
    main()
