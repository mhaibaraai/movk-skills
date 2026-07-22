#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0"]
# ///
"""大纲 JSON -> .pptx：克隆模板中已设计好的页面，再把文本回填进槽位。

用法：
    uv run scripts/build_pptx.py --outline outline.json --out 输出.pptx

不做程序化绘制——每一页都来自模板原页的深拷贝，装饰形状、配图、配色、字体全部保留，
产物与模板视觉一致，且是标准 OOXML，可在 PowerPoint/WPS 二次编辑。
大纲 schema 见 references/outline-schema.md。
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_util import (  # noqa: E402
    clone_slide,
    drop_slides,
    iter_text_shapes,
    load_index,
    resolve_template,
    set_text,
    source_path,
)


def slot_texts(page: dict, slots: list[dict], kind: str, section_no: int) -> dict[int, str]:
    """把大纲一页的内容摊到该模板页的槽位上：sid -> 文本。

    要点按槽位的 idx 取用（索引阶段已按几何邻近把标题与正文配成同一条）；decor 保持原样不动。
    no 只在章节页改写为章节序号——内容页上的 01/02/03 是版式自带的条目编号，动了就乱。
    """
    items = page.get("items", [])
    filled: dict[int, str] = {}

    for slot in slots:
        role, sid = slot["role"], slot["sid"]
        if role == "decor":
            continue
        if role == "no":
            if kind == "section" and section_no:
                filled[sid] = f"{section_no:02d}"
            continue
        if role == "title":
            filled[sid] = page.get("title", "")
        elif role == "subtitle":
            filled[sid] = page.get("subtitle", "")
        elif role in ("item_head", "item_body"):
            i = slot.get("idx", 0)
            item = items[i] if i < len(items) else {}
            filled[sid] = item.get("head" if role == "item_head" else "body", "")
    # 空串也要写回：宁可留空框，也不能把「单击此处添加文本」带进成品
    return filled


def render(outline: dict, out: Path, index_path: str | None = None, source: str | None = None) -> int:
    if index_path:  # 自带模板：显式索引 + 源 pptx，旁路 registry
        slides = json.loads(Path(index_path).read_text(encoding="utf-8"))["slides"]
        src_pptx = Path(source)
    else:
        key, _ = resolve_template(outline["template"])
        slides = load_index(key)["slides"]
        src_pptx = source_path(key)
    index = {s["i"]: s for s in slides}
    prs = Presentation(str(src_pptx))
    original = len(prs.slides)

    section_no = 0
    for page in outline["pages"]:
        src = page["src"]
        if src not in index:
            raise SystemExit(f"模板 {key} 没有第 {src} 页，检查大纲的 src 字段")
        kind = index[src]["kind"]
        if kind == "section":
            section_no += 1
        slide = clone_slide(prs, src)
        texts = slot_texts(page, index[src]["slots"], kind, section_no)
        for shape in iter_text_shapes(slide.shapes):
            if shape.shape_id in texts:
                set_text(shape, texts[shape.shape_id])

    drop_slides(prs, original)
    prs.save(str(out))
    return len(outline["pages"])


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outline", required=True)
    parser.add_argument("--out", default=None)
    parser.add_argument("--index", help="自带模板：显式索引 JSON 路径，旁路 registry")
    parser.add_argument("--source", help="自带模板：源 pptx 路径（与 --index 配套）")
    args = parser.parse_args()

    if bool(args.index) != bool(args.source):
        raise SystemExit("--index 与 --source 必须同时提供（自带模板）")
    outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
    out = Path(args.out or f"{outline.get('title', 'output')}.pptx")
    pages = render(outline, out, args.index, args.source)
    print(f"已生成 {out}（{pages} 页，模板：{outline['template']}）")


if __name__ == "__main__":
    main()
