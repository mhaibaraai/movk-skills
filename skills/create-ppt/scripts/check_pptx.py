#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0", "defusedxml>=0.7"]
# ///
"""成品自检：OOXML 合法性、残留占位与文本溢出。渲染后跑一遍，有问题就修好再重渲。

另会列出含待补槽（`__`）的页码——那是用户确认过留空的数据位，不是问题，只提示不阻断。

用法：
    uv run scripts/check_pptx.py --outline outline.json 输出.pptx

溢出判据取自模板索引里每个槽位的 cap（几何容量与模板原文案长度的较大值），
不重新猜几何——单靠几何会把 LOGO、目录这类模板自带的设计文字误报成溢出。
合法性检查覆盖 PowerPoint 会判「内容有问题」的几类结构错误——文本没问题但打不开
的成品同样是废片，不能只查文案。不依赖任何渲染器，任何环境都能跑。
"""
from __future__ import annotations

import argparse
import collections
import json
import posixpath
import re
import sys
import zipfile
from pathlib import Path

# 待检 pptx 可能来自用户上传，stdlib 解析器扛不住实体炸弹
from defusedxml import ElementTree as ET
from defusedxml.common import DefusedXmlException

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_util import iter_text_shapes, load_index  # noqa: E402

PLACEHOLDER_HINTS = ("单击此处", "单击添加", "此处添加", "本模板内", "请添加", "请自行替换", "XXXX", "{")
# 待补槽：连写的下划线算一处，`____` 不该报成两处
TODO_RE = re.compile(r"_{2,}")
TOLERANCE = 1.15  # 允许超 15%：cap 是估算值，卡太死会报一堆假警
SLIDE_RE = re.compile(r"ppt/slides/slide\d+\.xml$")
CNVPR_ID_RE = re.compile(r"<p:cNvPr[^>]*\bid=\"(\d+)\"")
RID_USE_RE = re.compile(r"r:(?:id|embed|link)=\"(rId\d+)\"")
RID_DECL_RE = re.compile(r"Id=\"(rId\d+)\"")
TARGET_RE = re.compile(r"Target=\"([^\"]+)\"(?:\s+TargetMode=\"(External)\")?")
# 图片可多页共用，模板自己就这么干；tags/chart 这类每形状私有的部件被共享就是损坏
SHAREABLE_PREFIX = ("ppt/media/",)
CHART_RE = re.compile(r"ppt/charts/chart\d+\.xml$")
EXTDATA_RID_RE = re.compile(r"externalData\s+r:id=\"(rId\d+)\"")
REL_LINE_RE = re.compile(r"Id=\"(rId\d+)\"[^>]*?Type=\"([^\"]+)\"")


def check_package(path: Path) -> list[str]:
    """OOXML 结构合法性：XML 可解析、页内 shape id 唯一、rId 不悬空、target 不缺失、
    每形状私有的部件（tags/chart）不跨页共享。这几类正是 PowerPoint 报「内容有问题」、
    修复时直接删掉相关形状的成因。"""
    issues: list[str] = []
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())
        for name in sorted(names):
            if name.endswith((".xml", ".rels")):
                try:
                    ET.fromstring(z.read(name))
                except DefusedXmlException as exc:
                    issues.append(f"  [XML 不安全] {name}：{exc}")
                except ET.ParseError as exc:
                    issues.append(f"  [XML 损坏] {name}：{exc}")

        for name in sorted(n for n in names if SLIDE_RE.match(n)):
            data = z.read(name).decode("utf-8", "ignore")
            dup = {i: c for i, c in collections.Counter(CNVPR_ID_RE.findall(data)).items() if c > 1}
            if dup:
                issues.append(f"  [重复 shape id] {name}：{dup}")
            rels_name = f"ppt/slides/_rels/{name.split('/')[-1]}.rels"
            declared = (
                set(RID_DECL_RE.findall(z.read(rels_name).decode("utf-8", "ignore")))
                if rels_name in names
                else set()
            )
            if missing := set(RID_USE_RE.findall(data)) - declared:
                issues.append(f"  [悬空 rId] {name}：{sorted(missing)}")

        refs: collections.Counter = collections.Counter()
        for name in sorted(n for n in names if n.endswith(".rels")):
            base = posixpath.dirname(posixpath.dirname(name))
            for target, external in TARGET_RE.findall(z.read(name).decode("utf-8", "ignore")):
                if external or target.startswith("http"):
                    continue
                full = posixpath.normpath(posixpath.join(base, target))
                if full not in names:
                    issues.append(f"  [target 缺失] {name} -> {target}")
                refs[full] += 1

        for part, count in sorted(refs.items()):
            if count > 1 and part.startswith("ppt/") and not part.startswith(SHAREABLE_PREFIX):
                if re.search(r"/(tags|charts|embeddings|diagrams)/", part):
                    issues.append(f"  [部件被 {count} 处共享] {part} — 该部件每形状私有，共享会被判损坏")

        # 图表 externalData 必须指向嵌入工作簿（package）：克隆重编号后若指到 chartStyle/
        # chartColorStyle，PowerPoint 会判「内容有问题」。rId 声明了但类型不对，悬空检查逮不到。
        for name in sorted(n for n in names if CHART_RE.match(n)):
            m = EXTDATA_RID_RE.search(z.read(name).decode("utf-8", "ignore"))
            if not m:
                continue
            rels_name = f"ppt/charts/_rels/{name.split('/')[-1]}.rels"
            types = dict(REL_LINE_RE.findall(z.read(rels_name).decode("utf-8", "ignore")))
            reltype = types.get(m.group(1), "")
            if not reltype.endswith("/package"):
                tail = reltype.rsplit("/", 1)[-1] or "缺失"
                issues.append(f"  [图表数据引用错位] {name}：externalData 指向 {tail}，应为嵌入工作簿")
    return issues


def check(outline: dict, path: Path) -> int:
    index = {s["i"]: s for s in load_index(outline["template"])["slides"]}
    prs = Presentation(str(path))
    issues: list[str] = check_package(path)
    todo: collections.Counter = collections.Counter()

    for n, (page, slide) in enumerate(zip(outline["pages"], prs.slides), 1):
        caps = {
            slot["sid"]: slot["cap"]
            for slot in index[page["src"]]["slots"]
            if "cap" in slot
        }
        for shape in iter_text_shapes(slide.shapes):
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if any(hint in text for hint in PLACEHOLDER_HINTS):
                issues.append(f"  [残留占位] 第 {n} 页：{text[:24]}")
            todo[n] += len(TODO_RE.findall(text))
            cap = caps.get(shape.shape_id)
            if cap and len(text) > cap * TOLERANCE:
                issues.append(f"  [文本溢出] 第 {n} 页：{len(text)} 字 / 上限 {cap} 字 — {text[:24]}")

    print(f"{path.name}：{len(prs.slides)} 页")
    print("\n".join(issues) if issues else "  通过：结构合法、无残留占位、无文本溢出")
    if pages := sorted(n for n, c in todo.items() if c):
        joined = "、".join(str(n) for n in pages)
        print(f"  待补数据：第 {joined} 页（共 {sum(todo.values())} 处 __ 槽位）— 提示，不阻断发布")
    return len(issues)  # 待补槽不计入：非零返回值会被 main 变成阻断


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outline", required=True)
    parser.add_argument("pptx")
    args = parser.parse_args()

    outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
    sys.exit(1 if check(outline, Path(args.pptx)) else 0)


if __name__ == "__main__":
    main()
