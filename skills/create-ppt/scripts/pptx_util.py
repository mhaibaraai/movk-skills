#!/usr/bin/env python3
"""模板复用的公共底座：文本槽遍历、幻灯片克隆、模板注册表加载。

被 index_template.py / build_pptx.py / make_outline.py 共用。
"""
from __future__ import annotations

import copy
import json
import re
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def load_registry() -> dict:
    return json.loads((ROOT / "templates" / "registry.json").read_text(encoding="utf-8"))


def resolve_template(name: str) -> tuple[str, dict]:
    """模板键或别名 → (标准键, 注册项)。命中不了抛错并列出可选值。"""
    reg = load_registry()
    if name in reg:
        return name, reg[name]
    for key, item in reg.items():
        if name in item.get("aliases", []):
            return key, item
    raise SystemExit(f"未知模板：{name}。可选：{'、'.join(reg)}")


def load_index(key: str) -> dict:
    _, item = resolve_template(key)
    return json.loads((ROOT / item["index"]).read_text(encoding="utf-8"))


def source_path(key: str) -> Path:
    _, item = resolve_template(key)
    return ROOT / item["source"]


def iter_text_shapes(shapes):
    """递归遍历（含组合）中所有带文本框的形状。"""
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_text_shapes(shape.shapes)
        elif shape.has_text_frame:
            yield shape


def shape_font_size(shape) -> float:
    """取形状内最大字号（pt）。全部继承自版式时返回 0。"""
    sizes = [
        run.font.size.pt
        for para in shape.text_frame.paragraphs
        for run in para.runs
        if run.font.size is not None
    ]
    return max(sizes) if sizes else 0.0


DEFAULT_FONT_PT = 18.0


def capacity(shape) -> int:
    """这个文本框能放下几个汉字：框宽高 ÷ 字号的几何估算，汉字按全角计。

    偏保守——形状常带自动缩放或允许居中溢出，所以索引阶段还要用模板自带的文案长度兜底
    （见 index_template.cap_of），单靠几何会把 LOGO、目录这类设计好的文字误判成溢出。
    """
    from pptx.util import Emu  # 延迟导入，让不碰 pptx 的调用方（make_outline）保持零依赖

    font = shape_font_size(shape) or DEFAULT_FONT_PT
    per_line = max(1, round(Emu(shape.width or 0).pt * 0.92 / font))
    lines = max(1, int(Emu(shape.height or 0).pt * 0.95 // (font * 1.25)))
    return per_line * lines


def set_text(shape, text: str) -> None:
    """回填文本并保留首个 run 的字体样式，多余 run/段落清除。"""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    runs = para.runs
    if not runs:
        para.add_run()
        runs = para.runs
    runs[0].text = text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)


PARTNAME_TMPL_RE = re.compile(r"^(.*?)(\d+)(\.[^.]+)$")
# 图片可以多页共用（模板自己就这么干），其余部件不行：tags 存的是 WPS 每形状私有的
# 智能图形元数据（KSO_WM_*），chart 各自挂着独立子部件。同一源页被排布复用时若共享
# 这些部件，PowerPoint 会判成「内容有问题」，修复时把相关形状整个删掉。
SHAREABLE = ("/image", "/video", "/audio", "/media")


def _relate_as(part, rid: str, target, reltype: str, is_external: bool = False) -> None:
    """在 part 上新建一条关系并沿用源 rId。

    不能用 relate_to：它按 _next_rId 重新发号，而部件 blob 内部的 r:id 引用
    （如 chart 的 <c:externalData r:id="rIdN">）是照抄过来的，重号后就指错部件，
    PowerPoint 会判「内容有问题」。新部件的 rels 是空的，沿用源 rId 不会冲突。
    """
    from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
    from pptx.opc.package import _Relationship

    mode = RTM.EXTERNAL if is_external else RTM.INTERNAL
    part.rels._rels[rid] = _Relationship(part.rels._base_uri, rid, reltype, mode, target)


def _dup_part(package, src_part, depth: int = 0):
    """把部件连同其子关系复制成一份独立的新部件——克隆页各用各的，不共享。"""
    from pptx.opc.package import Part

    match = PARTNAME_TMPL_RE.match(str(src_part.partname))
    if not match or depth > 4:  # partname 不带序号或层级过深，退回共享
        return src_part
    tmpl = f"{match.group(1)}%d{match.group(3)}"
    new = Part(package.next_partname(tmpl), src_part.content_type, package, src_part.blob)
    for rid, rel in src_part.rels.items():
        if rel.is_external:
            _relate_as(new, rid, rel.target_ref, rel.reltype, is_external=True)
        elif any(tag in rel.reltype for tag in SHAREABLE):
            _relate_as(new, rid, rel.target_part, rel.reltype)
        else:
            _relate_as(new, rid, _dup_part(package, rel.target_part, depth + 1), rel.reltype)
    return new


def clone_slide(prs, src_index: int):
    """把 prs 内第 src_index 页深拷贝为一张新页（同 package，图片/主题自动复用）。

    rId 必须重映射：沿用源 rId 会与新页自带的 layout 关系冲突，save() 时崩在 target_ref。
    """
    src = prs.slides[src_index]
    new = prs.slides.add_slide(src.slide_layout)
    for shape in list(new.shapes):  # 清掉版式带来的空占位符
        shape._element.getparent().remove(shape._element)

    rid_map: dict[str, str] = {}
    for element in src.shapes._spTree:
        if element.tag.endswith(("}nvGrpSpPr", "}grpSpPr")):
            continue
        copied = copy.deepcopy(element)
        for node in copied.iter():
            for attr, old_rid in list(node.attrib.items()):
                if not attr.startswith(REL_NS):
                    continue
                if old_rid not in rid_map:
                    rel = src.part.rels[old_rid]
                    if rel.is_external:
                        new_rid = new.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                    elif any(tag in rel.reltype for tag in SHAREABLE):
                        new_rid = new.part.relate_to(rel.target_part, rel.reltype)
                    else:
                        target = _dup_part(new.part.package, rel.target_part)
                        new_rid = new.part.relate_to(target, rel.reltype)
                    rid_map[old_rid] = new_rid
                node.set(attr, rid_map[old_rid])
        new.shapes._spTree.append(copied)
    return new


def drop_slides(prs, count: int) -> None:
    """删除最前面 count 张幻灯片（克隆完成后用于清掉模板原始页）。"""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst)[:count]:
        rid = sld_id.get(REL_NS + "id")
        if rid:
            prs.part.drop_rel(rid)
        sld_id_lst.remove(sld_id)
