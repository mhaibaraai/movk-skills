#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0"]
# ///
"""检查 .pptx 模板：列出母版、版式、占位符、主题色与字体。

用法：
    uv run scripts/inspect_template.py ppt/中石化PPT模板-1.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def fmt_emu(v) -> str:
    if v is None:
        return "?"
    return f"{Emu(v).inches:.2f}in"


def inspect(path: str) -> None:
    prs = Presentation(path)
    print(f"# {path}")
    print(f"slide size: {fmt_emu(prs.slide_width)} x {fmt_emu(prs.slide_height)}")
    print(f"slides: {len(prs.slides)}  layouts: {sum(len(m.slide_layouts) for m in prs.slide_masters)}")

    for mi, master in enumerate(prs.slide_masters):
        print(f"\n## master[{mi}]")
        for li, layout in enumerate(master.slide_layouts):
            phs = [
                f"{p.placeholder_format.idx}:{p.placeholder_format.type}:{p.name}"
                for p in layout.placeholders
            ]
            print(f"  layout[{li}] {layout.name!r} -> {phs}")

    print("\n## sample slides")
    for si, slide in enumerate(prs.slides):
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text.strip().replace("\n", " | ")[:40])
        print(f"  slide[{si}] layout={slide.slide_layout.name!r} :: {texts[:4]}")


if __name__ == "__main__":
    inspect(sys.argv[1] if len(sys.argv) > 1 else "ppt/中石化PPT模板-1.pptx")
