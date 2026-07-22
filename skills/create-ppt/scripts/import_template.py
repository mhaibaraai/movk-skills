#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0"]
# ///
"""把用户上传的 .pptx 取来并索引成临时模板，供 build_pptx/make_outline 的 --index/--source 旁路使用。

取数支持三种来源：完整 URL、平台相对路径（配 --base 拼域名）、本地文件。下载的只是用户
自己上传的这一个已知文件，不是联网检索数据——与「缺数据不联网」的规范不冲突。

用法：
    # 平台相对路径 + 全局 base（{{ 开始.other }} 给的是 /api/file/xxx）
    uv run skills/create-ppt/scripts/import_template.py /api/file/xxx --base https://llm.geosophon.com --out-dir tpl
    # 完整 URL
    uv run skills/create-ppt/scripts/import_template.py https://host/api/file/xxx --out-dir tpl
    # 本地文件
    uv run skills/create-ppt/scripts/import_template.py 我的模板.pptx --out-dir tpl

产出 tpl/source.pptx 与 tpl/index.json（含 template/source/sections/slides），并打印版式清单摘要。
下游（沙箱每轮全新，渲染轮需按同一 URL 重跑本脚本再消费）：
    make_outline.py --index tpl/index.json --pages N --title ... --sections ...  # 标签取自 index，--template 可覆盖
    build_pptx.py  --outline outline.json --index tpl/index.json --source tpl/source.pptx
    check_pptx.py  --outline outline.json --index tpl/index.json out.pptx
"""
from __future__ import annotations

import argparse
import collections
import json
import shutil
import sys
import urllib.request
import zipfile
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
from index_template import index_slide  # noqa: E402

MAX_BYTES = 100 * 1024 * 1024  # 下载体积上限，防异常/超大文件拖垮沙箱
TIMEOUT = 60
KIND_ORDER = ("cover", "toc", "section", "content", "table", "closing")


def acquire(source: str, base: str | None, dest: Path) -> None:
    """把源（完整 URL / 平台相对路径+base / 本地文件）落到 dest。"""
    url = None
    if source.startswith(("http://", "https://")):
        url = source
    elif base:
        url = base.rstrip("/") + "/" + source.lstrip("/")

    if url:
        req = urllib.request.Request(url, headers={"User-Agent": "create-ppt/import-template"})
        try:
            with urllib.request.urlopen(req, timeout=TIMEOUT) as resp:
                data = resp.read(MAX_BYTES + 1)
        except OSError as exc:
            raise SystemExit(f"下载模板失败：{url} —— {exc}")
        if len(data) > MAX_BYTES:
            raise SystemExit(f"下载体积超过上限 {MAX_BYTES // 1024 // 1024}MB，拒绝处理")
        dest.write_bytes(data)
        return

    src = Path(source)
    if not src.exists():
        raise SystemExit(f"本地文件不存在：{source}（要下载请给完整 URL 或配 --base）")
    shutil.copy2(src, dest)


def verify_pptx(path: Path) -> None:
    """确认是合法 pptx 包（含 ppt/presentation.xml），否则明确报错——上传/下载常拿到错内容。"""
    if not zipfile.is_zipfile(path):
        raise SystemExit("取到的文件不是 zip/pptx，无法作模板（可能下载到错误内容或非 pptx）")
    with zipfile.ZipFile(path) as z:
        if "ppt/presentation.xml" not in z.namelist():
            raise SystemExit("不是有效的 PowerPoint 包（缺 ppt/presentation.xml）")


def derive_sections(slides: list[dict]) -> list[str]:
    """从章节页/目录页样本尽力提取默认章节名——模型可用 --sections 覆盖。"""
    names = [
        slot["sample"]
        for s in slides
        if s["kind"] == "section"
        for slot in s["slots"]
        if slot["role"] == "title"
    ]
    if not names:
        toc = next((s for s in slides if s["kind"] == "toc"), None)
        if toc:
            names = [slot["sample"] for slot in toc["slots"] if slot["role"] == "item_head"]
    return list(dict.fromkeys(n for n in names if n))


def summary(name: str, slides: list[dict]) -> str:
    kinds = collections.Counter(s["kind"] for s in slides)
    content_items = sorted({s["items"] for s in slides if s["kind"] == "content" and s.get("items")})
    inv = " ".join(f"{k}{kinds[k]}" for k in KIND_ORDER if kinds.get(k))
    return (
        f"模板「{name}」已索引：共 {len(slides)} 页版式（{inv}）；"
        f"内容页可容纳要点数：{content_items or '—'}。"
        f"自动识别为草稿，封面/目录/章节页若有误判可在预览缩略图中核对。"
    )


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", help="完整 URL / 平台相对路径（配 --base）/ 本地 pptx")
    parser.add_argument("--base", help="平台域名，前缀到相对路径上（如 https://llm.geosophon.com）")
    parser.add_argument("--out-dir", default="tpl", help="产物目录（默认 tpl/）")
    parser.add_argument("--name", help="模板标签，默认「自带模板」；跨轮须保持一致")
    args = parser.parse_args()

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    src_pptx = out_dir / "source.pptx"

    acquire(args.source, args.base, src_pptx)
    verify_pptx(src_pptx)

    name = args.name or "自带模板"
    prs = Presentation(str(src_pptx))
    total = len(prs.slides)
    if total == 0:
        raise SystemExit("模板一页都没有，无法作为版式来源")
    slides = [index_slide(s, i, total) for i, s in enumerate(prs.slides)]
    index = {
        "template": name,
        "source": str(src_pptx),
        "sections": derive_sections(slides),
        "slides": slides,
    }
    (out_dir / "index.json").write_text(
        json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    print(summary(name, slides))
    print(f"index: {out_dir / 'index.json'}")
    print(f"source: {src_pptx}")


if __name__ == "__main__":
    main()
