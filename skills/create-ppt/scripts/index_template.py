#!/usr/bin/env python3
# /// script
# dependencies = ["python-pptx>=1.0"]
# ///
"""把一个设计好的 .pptx 索引成「页型 + 文本槽位」，供 build_pptx.py 克隆回填。

用法：
    uv run scripts/index_template.py assets/employee-safety-training.pptx > templates/index/employee-safety-training.json

输出为草稿：kind / role 由启发式推断，需人工或模型复核后落盘。
槽位一律按 shape id 定位——模板里 shape name 大量重复，不可作为标识。
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_util import capacity, iter_text_shapes, shape_font_size  # noqa: E402

DECOR_PATTERNS = [
    r"^\d{1,2}$",
    r"^(PART|Part)\s*\.?\s*\d+$",
    r"^[A-Za-z][A-Za-z\s\.&/-]*$",  # 纯英文装饰字（PROJECT APPLICATION、LOGO、CONTENTS）
    r"^第?[一二三四五六七八九十]+$",
]
TOC_HINTS = ("目录", "CONTENTS", "CONTENT")
CLOSING_HINTS = ("谢谢", "感谢", "THANKS", "THANK YOU")
SECTION_NO = re.compile(r"^(0?\d|PART\s*\.?\s*\d+)$", re.I)


def flatten(text: str) -> str:
    """去掉换行与空格再做匹配——模板里「目 录」「目\\n录」写法不一。"""
    return re.sub(r"\s+", "", text)


def is_decor(text: str) -> bool:
    flat = flatten(text)
    if any(hint in flat.upper() for hint in TOC_HINTS):
        return True
    return any(re.match(p, flat) for p in DECOR_PATTERNS)


def prominence(shape) -> float:
    """显著度：优先字号；模板里字号常继承自版式而取不到，退回形状高度（英寸×10）。"""
    size = shape_font_size(shape)
    return size or Emu(shape.height or 0).inches * 10


FINE_PRINT_PT = 9.0  # 小于此字号的文本是图注、细则、水印，不参与回填


def assign_roles(shapes: list, kind: str) -> list[str]:
    """按页型分派角色：装饰先剔除，其余按显著度排序后套用该页型的角色序列。"""
    roles: list[str | None] = []
    for shape in shapes:
        text = shape.text_frame.text.strip()
        font = shape_font_size(shape)
        if is_decor(text) or (0 < font < FINE_PRINT_PT):
            roles.append("no" if SECTION_NO.match(text) else "decor")
        else:
            roles.append(None)

    ranked = sorted(
        (i for i, role in enumerate(roles) if role is None),
        key=lambda i: prominence(shapes[i]),
        reverse=True,
    )
    for rank, i in enumerate(ranked):
        text = shapes[i].text_frame.text.strip()
        if kind in ("cover", "section", "closing"):
            roles[i] = ("title", "subtitle")[min(rank, 1)] if rank < 2 else "item_body"
        elif kind == "toc":
            roles[i] = "item_head"
        else:
            top_area = (shapes[i].top or 0) < int(2.2 * 914400)
            if rank == 0 and top_area and len(text) <= 20:
                roles[i] = "title"
            else:
                roles[i] = "item_head" if len(text) <= 12 else "item_body"
    return roles


def guess_kind(index: int, total: int, texts: list[str], has_table: bool, layout: str) -> str:
    joined = flatten("".join(texts)).upper()
    if index == 0:
        return "cover"
    if any(hint in joined for hint in TOC_HINTS):
        return "toc"
    if any(hint in joined for hint in CLOSING_HINTS):
        return "closing"
    if has_table:
        return "table"
    body = [t for t in texts if not is_decor(t)]
    has_no = any(SECTION_NO.match(flatten(t)) for t in texts)
    if has_no and (len(body) <= 2 or "节标题" in layout):
        return "section"
    if index == total - 1:
        return "closing"
    return "content"


CAP_MAX = {"title": 24, "subtitle": 34, "item_head": 14, "item_body": 80}


def cap_of(shape, role: str) -> int:
    """槽位字数上限。

    取几何容量与模板原文案长度的较大值——模板里已经排好的字必然放得下，
    是几何估算的下界；再按角色收口，避免正文槽因为框大就允许写成一整段。
    """
    fits = max(capacity(shape), len(shape.text_frame.text.strip()))
    return min(fits, CAP_MAX.get(role, 40))


def center(shape) -> tuple[float, float]:
    return (
        (shape.left or 0) + (shape.width or 0) / 2,
        (shape.top or 0) + (shape.height or 0) / 2,
    )


def pair_items(shapes: list, roles: list[str]) -> dict[int, int]:
    """把每个正文槽配到空间上最近的标题槽，得到「第几条要点」的 idx。

    标题数与正文数常不相等（版式里混有装饰性标题），各自按位置排序独立填充会错位——
    必须按几何邻近配对，才能保证一条要点的标题和正文落在同一个视觉分组里。
    """
    heads = [i for i, r in enumerate(roles) if r == "item_head"]
    bodies = [i for i, r in enumerate(roles) if r == "item_body"]
    idx_of = {head: n for n, head in enumerate(heads)}

    free = set(heads)
    next_idx = len(heads)
    for body in bodies:
        bx, by = center(shapes[body])
        nearest = min(
            free,
            key=lambda h: (center(shapes[h])[0] - bx) ** 2 + (center(shapes[h])[1] - by) ** 2,
            default=None,
        )
        if nearest is None:
            idx_of[body] = next_idx
            next_idx += 1
        else:
            idx_of[body] = idx_of[nearest]
            free.discard(nearest)
    return idx_of


ROW_BAND = int(0.6 * 914400)  # 0.6 英寸：同一横排内的高低差容忍度


def reading_order(shape) -> tuple[int, int]:
    """阅读顺序：先把 top 量化成行带再按 left 排。

    版式里同一排的卡片常有几毫米高低差（弧形排布、错落设计），
    直接按 top 精确排序会把「左中右」打乱成「中左右」，要点顺序就错了。
    """
    return ((shape.top or 0) // ROW_BAND, shape.left or 0)


def index_slide(slide, index: int, total: int) -> dict:
    shapes = [s for s in iter_text_shapes(slide.shapes) if s.text_frame.text.strip()]
    shapes.sort(key=reading_order)
    texts = [s.text_frame.text.strip() for s in shapes]
    has_table = any(s.has_table for s in slide.shapes)
    kind = guess_kind(index, total, texts, has_table, slide.slide_layout.name)

    roles = assign_roles(shapes, kind)
    idx_of = pair_items(shapes, roles)
    slots = []
    for i, (shape, role) in enumerate(zip(shapes, roles)):
        text = shape.text_frame.text.strip()
        slot = {"sid": shape.shape_id, "role": role, "sample": text[:24]}
        if role not in ("decor", "no"):
            slot["cap"] = cap_of(shape, role)
        if i in idx_of:
            slot["idx"] = idx_of[i]
        slots.append(slot)

    entry = {"i": index, "kind": kind, "slots": slots}
    # 封面/封底也可能有可填槽位（汇报人、汇报单位、日期），一并计入，否则会被清空成空框
    if idx_of and kind in ("content", "toc", "cover", "closing"):
        entry["items"] = max(idx_of.values()) + 1
    return entry


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--name", default=None, help="模板键，默认取文件名")
    args = parser.parse_args()

    path = Path(args.pptx)
    prs = Presentation(str(path))
    total = len(prs.slides)
    data = {
        "template": args.name or path.stem,
        "source": f"assets/{path.name}",
        "slides": [index_slide(s, i, total) for i, s in enumerate(prs.slides)],
    }
    print(json.dumps(data, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
